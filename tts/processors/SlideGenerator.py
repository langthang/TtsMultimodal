import os
from typing import Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from lxml import etree
from pptx.util import Inches
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Cm
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls
from AppConfig import AppConfig

class SlideGenerator:
    def __init__(self):
        self.slide_width = Inches(13.33)
        self.slide_height = Inches(7.5)
        self.config = AppConfig()

    def create_slide(self, title: str, content: str, background_image: str, output_file: str, translated_content: Optional[str] = None) -> str:
        """Create a slide with the given content and return the file path."""
        presentation = Presentation()
        presentation.slide_width = self.slide_width
        presentation.slide_height = self.slide_height

        # Create a blank slide
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])  # Layout 6 is blank
        
        # Add background if exists
        if background_image and os.path.exists(background_image):
            self._add_background(slide, background_image)

        # Add content
        if content:
            self._add_content(slide, content)
        if translated_content and self.config.activate_translation:
            self._add_translated_content(slide, translated_content)
        if title and self.config.enable_slide_title:
            self._add_title(slide, title)

        # Add logo to the top right
        self._add_logo(slide)

        # Save presentation
        presentation.save(output_file)
        return output_file

    def _add_background(self, slide, background_image: str):
        """Add background image to slide."""
        background_shape = slide.shapes.add_picture(
            background_image, 0, 0, 
            width=self.slide_width, 
            height=self.slide_height
        )

        slide.shapes._spTree.remove(background_shape._element)
        slide.shapes._spTree.insert(2, background_shape._element)

    def _add_title(self, slide, title: str):
        """Add title to slide."""
        # Create a new shape for title with initial dimensions
        title_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,  # Using regular rectangle
            Inches(0),  # Left position
            self.slide_height - Inches(4),  # Position from bottom
            Inches(0.1),  # Initial width, will be adjusted
            Inches(0.1)  # Initial height, will be adjusted
        )
        
        # Set transparent border
        title_shape.line.fill.solid()
        title_shape.line.fill.fore_color.rgb = RGBColor(*self.config.slide_text_background_color)
        title_shape.line.width = Pt(1)  # 1 point width
        if hasattr(title_shape.line.fill._xPr, 'solidFill'):
            srgbClr = title_shape.line.fill._xPr.solidFill.srgbClr
            self.SubElement(srgbClr, 'a:alpha', val='0')  # 0% opacity (completely transparent)
        
        # Set text frame properties
        text_frame = title_shape.text_frame
        text_frame.clear()
        text_frame.word_wrap = False  # Disable word wrap to get actual text width
        text_frame.margin_left = Inches(0.5)
        text_frame.margin_right = Inches(0.5)
        text_frame.margin_top = Inches(0.2)
        text_frame.margin_bottom = Inches(0.2)

        # Add title text
        p = text_frame.add_paragraph()
        p.text = title
        p.font.name = self.config.slide_title_font_name
        p.font.size = Pt(self.config.slide_title_font_size)
        p.font.color.rgb = RGBColor(*self.config.slide_title_font_color)
        p.alignment = PP_ALIGN.LEFT

        # Calculate approximate dimensions based on font size and text
        font_size_in_inches = self.config.slide_title_font_size / 72  # Convert points to inches
        num_lines = len(title.split('\n'))
        text_height = font_size_in_inches * num_lines
        
        # Calculate width based on the longest line
        longest_line = max(title.split('\n'), key=len)
        # Approximate width: each character is roughly 0.6 times the font size
        char_width = font_size_in_inches * 0.6
        text_width = len(longest_line) * char_width
        
        # Set shape dimensions with padding
        title_shape.height = Inches(text_height) + Inches(0.4)  # Add padding for margins
        title_shape.width = Inches(text_width) + Inches(1.0)  # Add padding for margins

        #title_shape.fill.solid()
        #title_shape.fill.fore_color.rgb = RGBColor(*self.config.slide_text_background_color)  # Background color
        self._set_shape_transparency_v3(title_shape, '0', angle='13440000', trans_angle='8100000')  # Fill: 0 degrees, Transparency: 135 degrees
    

    def _add_content(self, slide, content: str):
        """Add content to slide."""
        # Create a new shape for content at the bottom
        content_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,  # Using regular rectangle
            Inches(0),  # Left position
            self.slide_height - Inches(3),  # Position from bottom
            self.slide_width,  # Full width
            Inches(3)  # Height
        )
        
        # Set transparent border
        content_shape.line.fill.solid()
        content_shape.line.fill.fore_color.rgb = RGBColor(*self.config.slide_text_background_color)
        content_shape.line.width = Pt(1)  # 1 point width
        if hasattr(content_shape.line.fill._xPr, 'solidFill'):
            srgbClr = content_shape.line.fill._xPr.solidFill.srgbClr
            self.SubElement(srgbClr, 'a:alpha', val='0')  # 0% opacity (completely transparent)
        
        # Set text frame properties
        text_frame = content_shape.text_frame
        text_frame.clear()
        text_frame.word_wrap = True
        text_frame.margin_left = Inches(0.5)
        text_frame.margin_right = Inches(0.5)
        text_frame.margin_top = Inches(0.2)
        text_frame.margin_bottom = Inches(0.2)
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE  # Center vertically

        #content_shape.fill.solid()
        #content_shape.fill.fore_color.rgb = RGBColor(*self.config.slide_text_background_color)  # Background color
        self._set_shape_transparency_bottom(content_shape, '0', angle='21600000', trans_angle='10800000')  # Fill: 0 degrees, Transparency: 180 degrees
        

        # Process content to ensure left-aligned line-by-line formatting
        content_lines = [line.strip() for line in content.split("\n") if line.strip()]

        for line in content_lines:
            p = text_frame.add_paragraph()
            p.text = line
            p.font.name = self.config.slide_content_font_name
            p.font.color.rgb = RGBColor(*self.config.slide_content_font_color)
            p.font.size = Pt(self.config.slide_content_font_size)
            p.alignment = PP_ALIGN.CENTER  # Center horizontally

            # Remove bullet points and align text properly
            p._pPr.insert(0, etree.Element("{http://schemas.openxmlformats.org/drawingml/2006/main}buNone"))

    def _add_translated_content(self, slide, translated_content: str):
        """Add translated content to the top of the slide, opposite to the main content, with gradient effect."""
        # Create a new shape for translated content at the top
        translated_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,  # Using regular rectangle
            Inches(0),  # Left position
            Inches(0),  # Top position (top of slide)
            self.slide_width,  # Full width
            Inches(0)  # Height
        )

        # Set transparent border
        translated_shape.line.fill.solid()
        translated_shape.line.fill.fore_color.rgb = RGBColor(*self.config.slide_text_background_color)
        translated_shape.line.width = Pt(1)  # 1 point width
        if hasattr(translated_shape.line.fill._xPr, 'solidFill'):
            srgbClr = translated_shape.line.fill._xPr.solidFill.srgbClr
            self.SubElement(srgbClr, 'a:alpha', val='0')  # 0% opacity (completely transparent)

        # Set text frame properties
        text_frame = translated_shape.text_frame
        text_frame.clear()
        text_frame.word_wrap = True
        text_frame.margin_left = Inches(0.5)
        text_frame.margin_right = Inches(0.5)
        text_frame.margin_top = Inches(0.2)
        text_frame.margin_bottom = Inches(0.2)
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE  # Center vertically

        # Apply gradient transparency effect (top style)
        self._set_shape_transparency_top(translated_shape, '0', angle='10800000', trans_angle='5400000')  # Fill: 0 degrees, Transparency: 90 degrees

        # Process translated content to ensure left-aligned line-by-line formatting
        content_lines = [line.strip() for line in translated_content.split("\n") if line.strip()]

        for line in content_lines:
            p = text_frame.add_paragraph()
            p.text = line
            p.font.name = self.config.slide_content_font_name
            p.font.color.rgb = RGBColor(*self.config.slide_content_font_color)
            p.font.size = Pt(self.config.slide_content_font_size)
            p.alignment = PP_ALIGN.CENTER  # Center horizontally

            # Remove bullet points and align text properly
            p._pPr.insert(0, etree.Element("{http://schemas.openxmlformats.org/drawingml/2006/main}buNone"))


    def _add_logo(self, slide):
        """Add a small logo to the top right of the slide."""
        logo_path = os.path.join(os.path.dirname(__file__), "..", "logo", "logo-nobg-2.png")
        logo_path = os.path.abspath(logo_path)
        if not os.path.exists(logo_path):
            print(f"Logo file not found: {logo_path}")
            return
        # Set logo size and margin
        logo_width = Inches(1.5)
        logo_height = Inches(1.5)
        margin = Inches(3)
        left = self.slide_width - logo_width - margin
        top = margin
        slide.shapes.add_picture(logo_path, left, top, width=logo_width, height=logo_height)


    def SubElement(self, parent, tagname, **kwargs):
        element = OxmlElement(tagname)
        element.attrib.update(kwargs)
        parent.append(element)
        return element
    
    def _set_shape_transparency_bottom(self, shape, alpha, angle='16200000', trans_angle='16200000'):

        sp = shape._element

        # Access or create the <p:spPr> element (shape properties)
        spPr = sp.find("./p:spPr", namespaces={'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'})
        if spPr is None:
            spPr = parse_xml('<p:spPr xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"/>')
            sp.append(spPr)

        # Remove any existing fill elements
        for tag in ('a:noFill', 'a:solidFill', 'a:gradFill', 'a:blipFill', 'a:pattFill'):
            element = spPr.find(f"./{tag}", namespaces={'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
            if element is not None:
                spPr.remove(element)

        # Remove shadow (disable effects)
        # Remove existing effectLst if present
        effectLst = spPr.find('./a:effectLst', namespaces={'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
        if effectLst is not None:
            spPr.remove(effectLst)

        # Add explicit outer shadow with 0% opacity (invisible shadow)
        no_shadow_xml = '''
        <a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
        <a:outerShdw blurRad="40000" dist="20000" dir="5400000" algn="ctr" rotWithShape="0">
            <a:srgbClr val="000000">
            <a:alpha val="0"/>
            </a:srgbClr>
        </a:outerShdw>
        </a:effectLst>
        '''
        spPr.append(parse_xml(no_shadow_xml))


        # Define and insert gradient XML
        gradient_xml = f'''
        <a:gradFill rotWithShape="0" {nsdecls('a')}>
        <a:gsLst>
            <a:gs pos="100000">
            <a:srgbClr val="060606">
                <a:alpha val="0"/>
            </a:srgbClr>
            </a:gs>
            <a:gs pos="0">
            <a:srgbClr val="000000">
                <a:alpha val="100000"/>
            </a:srgbClr>
            </a:gs>
        </a:gsLst>
        <a:lin ang="16200000" scaled="1"/>
        </a:gradFill>
        '''
        spPr.append(parse_xml(gradient_xml))

    def _set_shape_transparency_top(self, shape, alpha, angle='16200000', trans_angle='16200000'):

        sp = shape._element

        # Access or create the <p:spPr> element (shape properties)
        spPr = sp.find("./p:spPr", namespaces={'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'})
        if spPr is None:
            spPr = parse_xml('<p:spPr xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"/>')
            sp.append(spPr)

        # Remove any existing fill elements
        for tag in ('a:noFill', 'a:solidFill', 'a:gradFill', 'a:blipFill', 'a:pattFill'):
            element = spPr.find(f"./{tag}", namespaces={'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
            if element is not None:
                spPr.remove(element)

        # Remove shadow (disable effects)
        # Remove existing effectLst if present
        effectLst = spPr.find('./a:effectLst', namespaces={'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
        if effectLst is not None:
            spPr.remove(effectLst)

        # Add explicit outer shadow with 0% opacity (invisible shadow)
        no_shadow_xml = '''
        <a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
        <a:outerShdw blurRad="40000" dist="20000" dir="5400000" algn="ctr" rotWithShape="0">
            <a:srgbClr val="000000">
            <a:alpha val="0"/>
            </a:srgbClr>
        </a:outerShdw>
        </a:effectLst>
        '''
        spPr.append(parse_xml(no_shadow_xml))


        # Define and insert gradient XML
        gradient_xml = f'''
        <a:gradFill rotWithShape="0" {nsdecls('a')}>
        <a:gsLst>
            <a:gs pos="0">
            <a:srgbClr val="060606">
                <a:alpha val="0"/>
            </a:srgbClr>
            </a:gs>
            <a:gs pos="100000">
            <a:srgbClr val="000000">
                <a:alpha val="100000"/>
            </a:srgbClr>
            </a:gs>
        </a:gsLst>
        <a:lin ang="16200000" scaled="1"/>
        </a:gradFill>
        '''
        spPr.append(parse_xml(gradient_xml))

    def _set_shape_transparency_v3(self, shape, alpha, angle='16200000', trans_angle='16200000'):
    
        sp = shape._element

        # Access or create the <p:spPr> element (shape properties)
        spPr = sp.find("./p:spPr", namespaces={'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'})
        if spPr is None:
            spPr = parse_xml('<p:spPr xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"/>')
            sp.append(spPr)

        # Remove any existing fill elements
        for tag in ('a:noFill', 'a:solidFill', 'a:gradFill', 'a:blipFill', 'a:pattFill'):
            element = spPr.find(f"./{tag}", namespaces={'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
            if element is not None:
                spPr.remove(element)

        # Remove shadow (disable effects)
        # Remove existing effectLst if present
        effectLst = spPr.find('./a:effectLst', namespaces={'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
        if effectLst is not None:
            spPr.remove(effectLst)

        # Add explicit outer shadow with 0% opacity (invisible shadow)
        no_shadow_xml = '''
        <a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
        <a:outerShdw blurRad="40000" dist="20000" dir="5400000" algn="ctr" rotWithShape="0">
            <a:srgbClr val="000000">
            <a:alpha val="0"/>
            </a:srgbClr>
        </a:outerShdw>
        </a:effectLst>
        '''
        spPr.append(parse_xml(no_shadow_xml))


        # Define and insert gradient XML
        gradient_xml = f'''
        <a:gradFill rotWithShape="0" {nsdecls('a')}>
        <a:gsLst>
            <a:gs pos="0">
            <a:srgbClr val="060606">
                <a:alpha val="0"/>
            </a:srgbClr>
            </a:gs>
            <a:gs pos="100000">
            <a:srgbClr val="000000">
                <a:alpha val="100000"/>
            </a:srgbClr>
            </a:gs>
        </a:gsLst>
        <a:lin ang="10800000" scaled="1"/>
        </a:gradFill>
        '''
        spPr.append(parse_xml(gradient_xml))
   
    def _set_shape_transparency(self, shape, alpha, angle='16200000', trans_angle='16200000'):
        """ Set the transparency (alpha) of a shape with gradient fill from light blue-green to transparent
        Args:
            shape: The shape to set transparency for
            alpha: The alpha value for transparency
            angle: The angle for the fill gradient in 60,000ths of a degree (default: 16200000 = 270 degrees)
            trans_angle: The angle for the transparency gradient in 60,000ths of a degree (default: 16200000 = 270 degrees)
        """

        sp = shape._element

        # Access or create the <p:spPr> element (shape properties)
        spPr = sp.find("./p:spPr", namespaces={'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'})
        if spPr is None:
            spPr = parse_xml('<p:spPr xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"/>')
            sp.append(spPr)

        # Remove any existing fill elements
        for tag in ('a:noFill', 'a:solidFill', 'a:gradFill', 'a:blipFill', 'a:pattFill'):
            element = spPr.find(f"./{tag}", namespaces={'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
            if element is not None:
                spPr.remove(element)

        # Remove shadow (disable effects)
        # Remove existing effectLst if present
        effectLst = spPr.find('./a:effectLst', namespaces={'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
        if effectLst is not None:
            spPr.remove(effectLst)

        # Add explicit outer shadow with 0% opacity (invisible shadow)
        no_shadow_xml = '''
        <a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
        <a:outerShdw blurRad="40000" dist="20000" dir="5400000" algn="ctr" rotWithShape="0">
            <a:srgbClr val="000000">
            <a:alpha val="0"/>
            </a:srgbClr>
        </a:outerShdw>
        </a:effectLst>
        '''
        spPr.append(parse_xml(no_shadow_xml))

        shape.fill.solid()

        # Remove existing fill
        if hasattr(shape.fill._xPr, 'solidFill'):
            shape.fill._xPr.remove(shape.fill._xPr.solidFill)
        
        # Create gradient fill
        gradFill = self.SubElement(shape.fill._xPr, 'a:gradFill')
        
        # Set gradient stops
        gsLst = self.SubElement(gradFill, 'a:gsLst')
        
        # First stop (0% position) - Light blue-green with full opacity
        gs1 = self.SubElement(gsLst, 'a:gs', pos='0')
        srgbClr1 = self.SubElement(gs1, 'a:srgbClr', val='ffffff')  # Light blue-green
        self.SubElement(srgbClr1, 'a:alpha', val=alpha)  # Full opacity
        
        # Second stop (100% position) - Light blue-green with specified transparency
        gs2 = self.SubElement(gsLst, 'a:gs', pos='100000')
        srgbClr2 = self.SubElement(gs2, 'a:srgbClr', val='a6d6d6')  # Light blue-green
        self.SubElement(srgbClr2, 'a:alpha', val=str('100000'))  # Specified transparency
        
        # Set linear gradient with specified angle
        lin = self.SubElement(gradFill, 'a:lin', ang=angle, scaled='1')
        
        # Set tile rectangle
        tileRect = self.SubElement(gradFill, 'a:tileRect')

   