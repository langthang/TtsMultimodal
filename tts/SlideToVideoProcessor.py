import os
from pptx import Presentation
from moviepy.editor import ImageClip


class SlideToVideoProcessor:
    def __init__(self, slide_file: str, video_length: int):
        """
        Initialize the processor with the slide file and video length.
        :param slide_file: Path to the PowerPoint slide file (.pptx).
        :param video_length: Length of the video in seconds.
        """
        self.slide_file = slide_file
        self.video_length = video_length

    def convert_slide_to_video(self, output_video_file: str):
        """
        Convert the slide to a video.
        :param output_video_file: Path to save the output video file (.mp4).
        """
        # Load the presentation
        presentation = Presentation(self.slide_file)

        # Extract the first slide
        slide = presentation.slides[0]

        # Save the slide as an image
        slide_image_path = os.path.splitext(self.slide_file)[0] + ".png"
        slide.shapes._parent.save(slide_image_path)

        # Create a video from the slide image
        clip = ImageClip(slide_image_path, duration=self.video_length)
        clip.write_videofile(output_video_file, fps=24)

        # Clean up the temporary slide image
        if os.path.exists(slide_image_path):
            os.remove(slide_image_path)

        print(f"Video saved to {output_video_file}")