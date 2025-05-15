import os
import json
import subprocess
from pdf2image import convert_from_path
from pptx.util import Inches
from mutagen.mp3 import MP3
from pptx import Presentation
from pptx.dml.color import RGBColor  # For setting text color
from pptx.enum.text import PP_ALIGN  # For text alignment
from pptx.util import Pt  # For setting font size
from lxml import etree  # For XML manipulation
from moviepy import *
from Conversations import Conversations
from GoogleTextToSpeech import GoogleTextToSpeech
from AppConfig import AppConfig

# Get the configuration instance
config = AppConfig()

# Mapping gender to Google TTS voice names
GENDER_TO_GOOGLE_TTS_VOICE_NAMES = {
    "female": ["en-US-Wavenet-F", "en-US-Wavenet-C", "en-US-Wavenet-E", "en-US-Wavenet-G", "en-US-Wavenet-H"],
    "male": ["en-US-Wavenet-A", "en-US-Wavenet-D", "en-US-Wavenet-B", "en-US-Wavenet-I", "en-US-Wavenet-J"],
    config.default_speaker: ["en-US-Wavenet-A"]
}

# Mapping language to Google TTS language codes
LANGUAGE_TO_GOOGLE_TTS_LANGUAGE_CODE = {
    "English": "en-US",
    "French": "fr-FR",
    "Spanish": "es-ES"
}


class TextToSpeechProcessor:
    def __init__(self, json_file: str):
        """Initialize the processor with the JSON file path."""
        self.json_file = json_file
        self.conversations_data = Conversations(json_file)
        self.google_tts = GoogleTextToSpeech()
        self.speaker_to_voice = {}
        self.used_voices = set()


    def assign_voices_to_speakers(self):
        """Assign unique voices to each speaker based on their gender."""
        for speaker_name, speaker in self.conversations_data.speakers.items():
            gender = speaker.gender.lower()
            available_voices = GENDER_TO_GOOGLE_TTS_VOICE_NAMES.get(gender, GENDER_TO_GOOGLE_TTS_VOICE_NAMES[config.default_speaker])

            # Find the next available voice
            for voice in available_voices:
                if voice not in self.used_voices:
                    self.speaker_to_voice[speaker_name] = voice
                    self.used_voices.add(voice)
                    break

        # If any speakers are left without a voice, assign them the default voice
        for speaker_name, speaker in self.conversations_data.speakers.items():
            if speaker_name not in self.speaker_to_voice:
                self.speaker_to_voice[speaker_name] = GENDER_TO_GOOGLE_TTS_VOICE_NAMES[config.default_speaker]


    def process_conversations(self):
        """Process each conversation line and synthesize speech."""
        # Get the directory of the JSON file
        json_dir = os.path.abspath(os.path.dirname(self.json_file))
        audios_dir = os.path.join(json_dir, "audio_conversations")
        slides_dir = os.path.join(json_dir, "slide_conversations")
        video_dir = os.path.join(json_dir, "video_conversations")

        # Create the audios directory if it doesn't exist
        if not os.path.exists(audios_dir):
            os.makedirs(audios_dir)

        # Create the slides directory if it doesn't exist
        if not os.path.exists(slides_dir):
            os.makedirs(slides_dir)

        # Create the videos directory if it doesn't exist
        if not os.path.exists(video_dir):
            os.makedirs(video_dir)
            
        language = self.conversations_data.language
        language_code = LANGUAGE_TO_GOOGLE_TTS_LANGUAGE_CODE.get(language, LANGUAGE_TO_GOOGLE_TTS_LANGUAGE_CODE[config.default_language])  # Default to English

        background_image = os.path.join(json_dir, "background.jpg")

        # Iterate through each conversation line
        for conversation in self.conversations_data.get_conversations():
            self.generate_speech_per_conversation(conversation, audios_dir, language_code)
            self.generate_slide_per_conversation(conversation, slides_dir, background_image)
            self.generate_video_per_conversation_slide(conversation, video_dir, conversation.slide, conversation.audio, conversation.audio_length)


    def generate_speech_per_conversation(self, conversation, audios_dir, language_code):
        """
        Process a single conversation line and synthesize speech.
        :param conversation: A single conversation object.
        :param audios_dir: Directory to save the audio file.
        :param language_code: Language code for TTS synthesis.
        """
        order = conversation.order
        text_to_speak = conversation.text
        speaker = conversation.speaker.name
        voice_name = self.speaker_to_voice.get(speaker, GENDER_TO_GOOGLE_TTS_VOICE_NAMES[config.default_speaker][0])  # Default to a neutral voice

        # Generate the output file path inside the audios directory
        audio_file = os.path.join(audios_dir, f"{order}_{speaker}.mp3")

        # Call synthesize_speech
        self.google_tts.synthesize_speech(
            text=text_to_speak,
            output_filename=audio_file,
            voice_name=voice_name,
            gender=conversation.speaker.gender.upper(),
            language_code=language_code
        )

        audio = MP3(audio_file)
        # Convert length to milliseconds
        audio_length = int(audio.info.length * 1000)  # Length in milliseconds

        conversation.audio_length = audio_length  # Update the conversation object with the audio length
        conversation.audio = audio_file  # Update the conversation object with the audio file path

        print(f"Synthesized speech for {order} : {speaker}: {text_to_speak} : {voice_name} : {audio_length} : {audio_file}")


    def generate_slide_per_conversation(self, conversation, slides_dir, background_image):
        """
        Generate a slide for a single conversation.
        :param conversation: A single conversation object.
        :param slides_dir: Directory to save the generated slide.
        :param background_image: Path to the background image.
        """
        order = conversation.order
        speaker = conversation.speaker.name
        text = conversation.text

        # Create a new presentation
        presentation = Presentation()
        presentation.slide_width = Inches(13.33)  # 16 inches
        presentation.slide_height = Inches(7.5)  # 9 inches

        # Add a slide with a title and content
        slide_layout = presentation.slide_layouts[1]  # Title and Content layout
        slide = presentation.slides.add_slide(slide_layout)

        # Set the background image for the slide
        if os.path.exists(background_image):
            # Add an image to the slide
            left = top = 0
            slide_width = presentation.slide_width
            slide_height = presentation.slide_height
            background_shape = slide.shapes.add_picture(
                background_image, left, top, width=slide_width, height=slide_height
            )
            # Move the background image to the back
            slide.shapes._spTree.remove(background_shape._element)
            slide.shapes._spTree.insert(2, background_shape._element)  # use the number that does the appropriate job

        # Set the title and content
        title = slide.shapes.title
        title.text = f"{speaker}"
        content = slide.placeholders[1]

        content.text_frame.paragraphs[0].text = text
        content.text_frame.paragraphs[0].font.size = Pt(22)
        content.text_frame.paragraphs[0].font.name = "Corbel"  # Set font to Corbel
        content.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # Set text color to white
        content.text_frame.paragraphs[0].font.bold = True
        content.text_frame.paragraphs[0]._pPr.insert(0, etree.Element("{http://schemas.openxmlformats.org/drawingml/2006/main}buNone"))
        content.text_frame.word_wrap = True  # Optional: Ensure text wraps within the placeholder
        content.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Set the font for the title as well
        for paragraph in title.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = "Corbel"  # Set font to Corbel
                run.font.color.rgb = RGBColor(255, 255, 255)  # Set text color to white
                run.font.size = Pt(24)  # Set font size to 20

        # Generate the output file path inside the slides directory
        slide_file = os.path.join(slides_dir, f"{order}_{speaker}.pptx")

        # Save the presentation
        presentation.save(slide_file)

        # Update the conversation object with the slide file path
        conversation.slide = slide_file

        print(f"Generated slide for {order}: {slide_file}")


    def generate_video_per_conversation_slide(self, conversation, video_dir, slide_file, audio_file, audio_length):
        """
        Convert the slide to a video with audio.
        :param video_dir: Directory to save the generated video.
        :param slide_file: Path to the slide file (.pptx).
        :param audio_file: Path to the audio file (.mp3).
        :param audio_length: Length of the video in seconds (equal to the audio length).
        """
        # Ensure the video directory exists
        if not os.path.exists(video_dir):
            os.makedirs(video_dir)

        # Generate the output video file path
        video_file = os.path.join(video_dir, os.path.splitext(os.path.basename(slide_file))[0] + ".avi")

        # Save the slide as an image
        slide_image_dir = os.path.abspath(os.path.dirname(slide_file))
        
        if config.slide_generation_mode_pdf:
            # Convert slide to pdf
            try:
                subprocess.run(
                    [
                        "soffice",
                        "--headless",
                        "--convert-to",
                        "pdf",
                        "--outdir",
                        slide_image_dir,
                        slide_file,
                    ],
                    check=True,
                )
            except subprocess.CalledProcessError as e:
                print(f"Error converting slide to pdf: {e}")
                return None
            
            slide_image_file_pdf = os.path.join(slide_image_dir, os.path.splitext(os.path.basename(slide_file))[0] + ".pdf")
            
            # Check if the image file exists
            if not os.path.exists(slide_image_file_pdf):
                print(f"PDF file not found: {slide_image_file_pdf}")
                return None
            
            slide_image_file = os.path.join(slide_image_dir, os.path.splitext(os.path.basename(slide_file))[0] + ".png")
               
            # Convert PDF to images (high DPI = better quality)
            images = convert_from_path(slide_image_file_pdf, dpi=900)
            images[0].save(slide_image_file, 'PNG')
        else:
            # Convert slide to png
            try:
                subprocess.run(
                    [
                        "soffice",
                        "--headless",
                        "--convert-to",
                        "png",
                        "--outdir",
                        slide_image_dir,
                        slide_file,
                    ],
                    check=True,
                )
            except subprocess.CalledProcessError as e:
                print(f"Error converting slide to image: {e}")
                return None

            slide_image_file = os.path.join(slide_image_dir, os.path.splitext(os.path.basename(slide_file))[0] + ".png")
            # To comment if generate PDF ================

        # Check if the image file exists
        if not os.path.exists(slide_image_file):
            print(f"Image file not found: {slide_image_file}")
            return None

        # Load the image file and create an ImageClip
        image_clip = ImageClip(img=slide_image_file, duration=(audio_length / 1000) + 1)  # Convert milliseconds to seconds
        image_clip.write_videofile(video_file, fps=60, codec="png", audio=audio_file, audio_codec="aac")

        # Clean up the temporary slide image
        # if os.path.exists(slide_image_file):
        #     os.remove(slide_image_file)

        conversation.video = video_file  # Update the conversation object with the video file path

        print(f"Generated video with audio: {video_file}")


    def merge_conversation_videos(self):
        """
        Merge all conversation videos into one video.
        :param conversations: List of Conversation objects.
        :param output_video_file: Path to save the merged video.
        """
        video_clips = []
        output_video_file = os.path.abspath(self.json_file).replace(".json", "conversations_merged_video.avi")

        # Collect video clips from conversations
        for conversation in self.conversations_data.get_conversations():
            if conversation.video and os.path.exists(conversation.video):
                video_clips.append(VideoFileClip(conversation.video))
            else:
                print(f"Video not found for conversation {conversation.order}: {conversation.video}")

        # Merge all video clips
        if video_clips:
            final_video = concatenate_videoclips(video_clips, method="compose")
            final_video.write_videofile(output_video_file, codec="png", audio_codec="aac", fps=60)
            print(f"Merged video saved to {output_video_file}")
        else:
            print("No videos found to merge.")


    def generate_speech_per_new_word(self, new_word, audios_dir, language_code):
        """
        Process each new word and synthesize speech.
        :param new_words: List of NewWord objects.
        :param audios_dir: Directory to save the audio files.
        :param language_code: Language code for TTS synthesis.
        """
        order = getattr(new_word, 'order', 0)  # Default to 0 if 'order' is missing
        word = getattr(new_word, 'word', None)
        meaning = getattr(new_word, 'meaning', None)
        example = getattr(new_word, 'example', None)
        voice_name = self.speaker_to_voice.get(config.default_speaker, GENDER_TO_GOOGLE_TTS_VOICE_NAMES[config.default_speaker][0])  # Default to a neutral voice


        # Combine the word, meaning, and example into a single text for speech synthesis
        if order == 0 or word is None or meaning is None:
            text_to_speak = f"{example}"
        else:
            text_to_speak = f"( … )( … ) {word}, ( … )( … ). Meaning ( … )( … ) {meaning} ( … ). Example ( … )( … )  {example}."

        # Generate the output file path inside the audios directory
        audio_file = os.path.join(audios_dir, f"new_word_{order}.mp3")

        # Call synthesize_speech
        self.google_tts.synthesize_speech(
            text=text_to_speak,
            output_filename=audio_file,
            voice_name=GENDER_TO_GOOGLE_TTS_VOICE_NAMES[config.default_speaker][0],
            gender=config.default_speaker,
            language_code=language_code
        )

        audio = MP3(audio_file)
        # Convert length to milliseconds
        audio_length = int(audio.info.length * 1000)  # Length in milliseconds

        new_word.audio = audio_file
        new_word.audio_length = audio_length

        print(f"Synthesized speech for new word {order} : {config.default_speaker}: {text_to_speak} : {voice_name} : {audio_length} : {audio_file}")

    
    def generate_slide_per_new_word(self, new_word, slides_dir, background_image):
        """
        Generate a slide for a single new word.
        :param new_word: A single NewWord object.
        :param slides_dir: Directory to save the generated slide.
        :param background_image: Path to the background image.
        """
        order = getattr(new_word, 'order', 0)  # Default to 0 if 'order' is missing
        word = getattr(new_word, 'word', None)
        meaning = getattr(new_word, 'meaning', None)
        example = getattr(new_word, 'example', None)

        # Create a new presentation
        presentation = Presentation()
        presentation.slide_width = Inches(13.33)  # 16 inches
        presentation.slide_height = Inches(7.5)  # 9 inches

        # Add a slide with a title and content
        slide_layout = presentation.slide_layouts[1]  # Title and Content layout
        slide = presentation.slides.add_slide(slide_layout)

        # Set the background image for the slide
        if os.path.exists(background_image):
            # Add an image to the slide
            left = top = 0
            slide_width = presentation.slide_width
            slide_height = presentation.slide_height
            background_shape = slide.shapes.add_picture(
                background_image, left, top, width=slide_width, height=slide_height
            )
            # Move the background image to the back
            slide.shapes._spTree.remove(background_shape._element)
            slide.shapes._spTree.insert(2, background_shape._element)  # use the number that does the appropriate job

        # Set the title and content
        title = slide.shapes.title
        if word != None:
            title.text = f"{word}"
        content = slide.placeholders[1]

        # if meaning != None:
        #     content.text_frame.paragraphs[0].text = f"\nMeaning: {meaning}\nExample: {example}"
        # content.text_frame.paragraphs[0].font.size = Pt(22)
        # content.text_frame.paragraphs[0].font.name = "Corbel"  # Set font to Corbel
        # content.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # Set text color to white
        # content.text_frame.paragraphs[0].font.bold = True
        # content.text_frame.paragraphs[0]._pPr.insert(0, etree.Element("{http://schemas.openxmlformats.org/drawingml/2006/main}buNone"))
        # content.text_frame.word_wrap = True  # Optional: Ensure text wraps within the placeholder
        # content.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

        content.text_frame.clear()  # Clear existing content

        # Add a paragraph for the meaning
        if meaning != None:
            meaning_paragraph = content.text_frame.add_paragraph()
            meaning_paragraph.text = "Meaning: "
            meaning_run = meaning_paragraph.add_run()
            meaning_run.text = meaning
            meaning_run.font.italic = True  # Make only the meaning italic
            meaning_paragraph.font.size = Pt(22)
            meaning_paragraph.font.name = "Corbel"
            meaning_paragraph.font.color.rgb = RGBColor(255, 255, 255)

            # Add a paragraph for the example
            example_paragraph = content.text_frame.add_paragraph()
            example_paragraph.text = "Example: "
            example_run = example_paragraph.add_run()
            example_run.text = example
            example_run.font.italic = True  # Make only the example italic
            example_paragraph.font.size = Pt(22)
            example_paragraph.font.name = "Corbel"
            example_paragraph.font.color.rgb = RGBColor(255, 255, 255)

        # Set the font for the title as well
        for paragraph in title.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = "Corbel"  # Set font to Corbel
                run.font.color.rgb = RGBColor(255, 255, 255)  # Set text color to white
                run.font.size = Pt(24)  # Set font size to 20

        # Generate the output file path inside the slides directory
        slide_file = os.path.join(slides_dir, f"new_word_{order}.pptx")

        # Save the presentation
        presentation.save(slide_file)

        # Update the new_word object with the slide file path
        new_word.slide = slide_file

        print(f"Generated slide for new word {order}: {slide_file}")


    def generate_video_per_new_word_slide(self, new_word, video_dir, slide_file, audio_file, audio_length):
        """
        Convert the slide for a new word to a video with audio.
        :param new_word: A single NewWord object.
        :param video_dir: Directory to save the generated video.
        :param slide_file: Path to the slide file (.pptx).
        :param audio_file: Path to the audio file (.mp3).
        :param audio_length: Length of the video in milliseconds (equal to the audio length).
        """
        # Ensure the video directory exists
        if not os.path.exists(video_dir):
            os.makedirs(video_dir)

        # Generate the output video file path
        video_file = os.path.join(video_dir, os.path.splitext(os.path.basename(slide_file))[0] + ".mp4")

        # Save the slide as an image
        slide_image_dir = os.path.abspath(os.path.dirname(slide_file))

        if config.slide_generation_mode_pdf:
            try:
                subprocess.run(
                    [
                        "soffice",
                        "--headless",
                        "--convert-to",
                        "pdf",
                        "--outdir",
                        slide_image_dir,
                        slide_file,
                    ],
                    check=True,
                )
            except subprocess.CalledProcessError as e:
                print(f"Error converting slide to pdf: {e}")
                return None
            
            slide_image_file_pdf = os.path.join(slide_image_dir, os.path.splitext(os.path.basename(slide_file))[0] + ".pdf")
            
            # Check if the image file exists
            if not os.path.exists(slide_image_file_pdf):
                print(f"PDF file not found: {slide_image_file_pdf}")
                return None
            
            slide_image_file = os.path.join(slide_image_dir, os.path.splitext(os.path.basename(slide_file))[0] + ".png")
            
            # Convert PDF to images (high DPI = better quality)
            images = convert_from_path(slide_image_file_pdf, dpi=900)
            images[0].save(slide_image_file, 'PNG')
        else:
            # Convert slide to png
            try:
                subprocess.run(
                    [
                        "soffice",
                        "--headless",
                        "--convert-to",
                        "png",
                        "--outdir",
                        slide_image_dir,
                        slide_file,
                    ],
                    check=True,
                )
            except subprocess.CalledProcessError as e:
                print(f"Error converting slide to image: {e}")
                return None

            slide_image_file = os.path.join(slide_image_dir, os.path.splitext(os.path.basename(slide_file))[0] + ".png")

        # Check if the image file exists
        if not os.path.exists(slide_image_file):
            print(f"Image file not found: {slide_image_file}")
            return None
        
        # Load the image file and create an ImageClip
        image_clip = ImageClip(img=slide_image_file, duration=(audio_length / 1000) + 1)  # Convert milliseconds to seconds
        image_clip.write_videofile(video_file, fps=60, codec="png", audio=audio_file, audio_codec="aac")

        # Clean up the temporary slide image
        # if os.path.exists(slide_image_file):
        #     os.remove(slide_image_file)

        # Update the new_word object with the video file path
        new_word.video = video_file

        print(f"Generated video for new word {new_word.order}: {video_file}")


    def process_new_words(self):
        """Process each new word and generate speech, slides, and videos."""
        # Get the directory of the JSON file
        json_dir = os.path.abspath(os.path.dirname(self.json_file))
        audios_dir = os.path.join(json_dir, "audio_new_words")
        slides_dir = os.path.join(json_dir, "slide_new_words")
        video_dir = os.path.join(json_dir, "video_new_words")

        # Create the directories if they don't exist
        if not os.path.exists(audios_dir):
            os.makedirs(audios_dir)
        if not os.path.exists(slides_dir):
            os.makedirs(slides_dir)
        if not os.path.exists(video_dir):
            os.makedirs(video_dir)

        language = self.conversations_data.language
        language_code = LANGUAGE_TO_GOOGLE_TTS_LANGUAGE_CODE.get(language, LANGUAGE_TO_GOOGLE_TTS_LANGUAGE_CODE[config.default_language])  # Default to English

        background_image = os.path.join(json_dir, "background.jpg")

        # Iterate through each new word
        for new_word in self.conversations_data.get_new_words():
            # Generate speech for the new word
            self.generate_speech_per_new_word(new_word, audios_dir, language_code)

            # Generate a slide for the new word
            self.generate_slide_per_new_word(new_word, slides_dir, background_image)

            # Generate a video for the new word
            self.generate_video_per_new_word_slide(new_word, video_dir, new_word.slide, new_word.audio, new_word.audio_length)


    def merge_new_word_videos(self):
        """
        Merge all new word videos into one video.
        """
        video_clips = []
        output_video_file = os.path.abspath(self.json_file).replace(".json", "_new_words_merged_video.mp4")

        # Collect video clips from new words
        for new_word in self.conversations_data.get_new_words():
            if new_word.video and os.path.exists(new_word.video):
                video_clips.append(VideoFileClip(new_word.video))
            else:
                print(f"Video not found for new word {new_word.order}: {new_word.video}")

        # Merge all video clips
        if video_clips:
            final_video = concatenate_videoclips(video_clips, method="compose")
            final_video.write_videofile(output_video_file, codec="libx264", audio_codec="aac", fps=24)
            print(f"Merged new words video saved to {output_video_file}")
        else:
            print("No videos found to merge for new words.")


    def merge_all_videos(self):
        """
        Merge all conversation and new word videos into one video.
        """
        video_clips = []
        output_video_file = os.path.abspath(self.json_file).replace(".json", "_merged_video.mp4")

        # Collect video clips from conversations
        for conversation in self.conversations_data.get_conversations():
            if conversation.video and os.path.exists(conversation.video):
                video_clips.append(VideoFileClip(conversation.video))
            else:
                print(f"Video not found for conversation {conversation.order}: {conversation.video}")

        # Collect video clips from new words
        for new_word in self.conversations_data.get_new_words():
            if new_word.video and os.path.exists(new_word.video):
                video_clips.append(VideoFileClip(new_word.video))
            else:
                print(f"Video not found for new word {new_word.order}: {new_word.video}")

        # Merge all video clips
        if video_clips:
            final_video = concatenate_videoclips(video_clips, method="compose")
            final_video.write_videofile(output_video_file, codec="libx264", audio_codec="aac", fps=24)
            print(f"Merged all videos saved to {output_video_file}")
        else:
            print("No videos found to merge.")


    def save_decorated_conversations_and_new_words(self):
        """
        Save the updated conversations to a new JSON file with the name json_file_decorated.
        """
        # Load the original JSON data
        with open(self.json_file, "r") as file:
            data = json.load(file)

        # Update the conversations in the JSON data
        data["conversations"] = [
            {
                "order": conversation.order,
                "speaker": conversation.speaker.name,
                "text": conversation.text,
                "slide": conversation.slide,
                "video": conversation.video,
                "audio_length": conversation.audio_length,
                "audio": conversation.audio
            }
            for conversation in self.conversations_data.get_conversations()
        ]

        # Update the new_words in the JSON data
        data["new_words"] = [
            {
                "order": new_word.order,
                "word": new_word.word,
                "meaning": new_word.meaning,
                "example": new_word.example,
                "slide": new_word.slide,
                "video": new_word.video,
                "audio_length": new_word.audio_length,
                "audio": new_word.audio
            }
            for new_word in self.conversations_data.get_new_words()
        ]

        # Generate the new file name
        decorated_file = os.path.splitext(self.json_file)[0] + "_decorated.json"

        # Save the updated JSON data to the new file
        with open(decorated_file, "w") as file:
            json.dump(data, file, indent=4)

        print(f"Updated conversations and new words saved to {decorated_file}")


    def run(self):
        """Run the entire text-to-speech processing pipeline."""
        self.assign_voices_to_speakers()
        self.process_conversations()
        self.merge_conversation_videos()
        self.process_new_words()
        self.merge_new_word_videos()
        self.merge_all_videos()
        self.save_decorated_conversations_and_new_words()