import os
from pptx import Presentation
from pptx.dml.color import RGBColor  # For setting text color
from pptx.enum.text import PP_ALIGN  # For text alignment
from pptx.util import Pt  # For setting font size
from lxml import etree  # For XML manipulation
from tts.Conversations import Conversations


class TextToSlideProcessor:
    def __init__(self, json_file: str):
        """Initialize the processor with the JSON file path."""
        self.json_file = json_file
        self.conversations_data = Conversations(json_file)

    def generate_slide_for_conversation(self, conversation, slides_dir, background_image):
        """
        Generate a slide for a single conversation.
        :param conversation: A single conversation object.
        :param slides_dir: Directory to save the generated slide.
        :param background_image: Path to the background image.
        """
        order = conversation.order
        speaker = conversation.speaker.name
        text = conversation.text

        # Create a new presentation
        presentation = Presentation()

        # Add a slide with a title and content
        slide_layout = presentation.slide_layouts[1]  # Title and Content layout
        slide = presentation.slides.add_slide(slide_layout)

        # Set the background image for the slide
        if os.path.exists(background_image):
            # Add an image to the slide
            left = top = 0
            slide_width = presentation.slide_width
            slide_height = presentation.slide_height
            background_shape = slide.shapes.add_picture(
                background_image, left, top, width=slide_width, height=slide_height
            )
            # Move the background image to the back
            slide.shapes._spTree.remove(background_shape._element)
            slide.shapes._spTree.insert(2, background_shape._element)  # use the number that does the appropriate job

        # Set the title and content
        title = slide.shapes.title
        title.text = f"{speaker}"
        content = slide.placeholders[1]

        content.text_frame.paragraphs[0].text = text
        content.text_frame.paragraphs[0].font.size = Pt(22)
        content.text_frame.paragraphs[0].font.name = "Corbel"  # Set font to Corbel
        content.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # Set text color to white
        content.text_frame.paragraphs[0].font.bold = True
        content.text_frame.paragraphs[0]._pPr.insert(0, etree.Element("{http://schemas.openxmlformats.org/drawingml/2006/main}buNone"))
        content.text_frame.word_wrap = True  # Optional: Ensure text wraps within the placeholder
        content.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Set the font for the title as well
        for paragraph in title.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = "Corbel"  # Set font to Corbel
                run.font.color.rgb = RGBColor(255, 255, 255)  # Set text color to white
                run.font.size = Pt(24)  # Set font size to 20

        # Generate the output file path inside the slides directory
        output_file_slide = os.path.join(slides_dir, f"{order}_{speaker}.pptx")

        # Save the presentation
        presentation.save(output_file_slide)
        print(f"Generated slide for {order}: {output_file_slide}")

    def process_conversations(self):
        """Generate one slide per conversation line."""
        # Get the directory of the JSON file
        json_dir = os.path.dirname(self.json_file)
        slides_dir = os.path.join(json_dir, "slide_conversations")
        background_image = os.path.join(os.path.abspath(json_dir), "background.jpg")
        print(f"Background image path: {background_image}")

        # Create the slides directory if it doesn't exist
        if not os.path.exists(slides_dir):
            os.makedirs(slides_dir)

        # Iterate through each conversation line
        for conversation in self.conversations_data.get_conversations():
            self.generate_slide_for_conversation(conversation, slides_dir, background_image)


    def run(self):
        """Run the entire text-to-slide processing pipeline."""
        self.process_conversations()